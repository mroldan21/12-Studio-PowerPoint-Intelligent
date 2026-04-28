"""
core/pptx_engine.py — Motor PPTX Mejorado
Funciones avanzadas de dibujo para temas visuales complejos.
"""

import io
from pathlib import Path
from typing import Optional, BinaryIO, Tuple, List, Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import requests


class PPTXEngine:
    """
    Motor completo para manipulación de presentaciones PowerPoint.
    """

    def __init__(self):
        self.prs: Optional[Presentation] = None

    # ── Lifecycle ──────────────────────────────────────────────────────────

    def create_blank(self, slide_count: int = 1) -> None:
        """Crea presentación nueva con slides en blanco."""
        self.prs = Presentation()
        blank_layout = self.prs.slide_layouts[6]  # blank layout
        for _ in range(slide_count):
            self.prs.slides.add_slide(blank_layout)

    def load(self, source: BinaryIO) -> None:
        """Carga PPTX existente."""
        self.prs = Presentation(source)

    def load_from_path(self, path: str) -> None:
        """Carga PPTX desde archivo."""
        p = Path(path)
        if not p.exists():
            raise FileNotFoundError(f"Template not found: {path}")
        self.prs = Presentation(str(p))

    def save(self, target: BinaryIO) -> None:
        """Guarda presentación."""
        if not self.prs:
            raise RuntimeError("No presentation loaded")
        self.prs.save(target)

    def load_template(self, template_path: str) -> None:
        """Carga una plantilla PPTX como base."""
        path = Path(template_path)
        if not path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}")
        self.prs = Presentation(str(path))

    def set_slide_dimensions(self, width_inches: float, height_inches: float) -> None:
        """Establece dimensiones de la diapositiva."""
        if not self.prs:
            raise RuntimeError("No presentation loaded")
        self.prs.slide_width = Inches(width_inches)
        self.prs.slide_height = Inches(height_inches)

    # ── Slide Management ───────────────────────────────────────────────────

    @property
    def slide_count(self) -> int:
        return len(self.prs.slides) if self.prs else 0

    def get_slide(self, index: int):
        if not self.prs:
            raise RuntimeError("No presentation loaded")
        if index >= len(self.prs.slides):
            raise IndexError(f"Slide {index} does not exist")
        return self.prs.slides[index]

    def add_blank_slide(self):
        """Añade slide en blanco."""
        layout = self.prs.slide_layouts[6]
        return self.prs.slides.add_slide(layout)

    # ── Notes ──────────────────────────────────────────────────────────────

    def set_notes(self, slide_index: int, text: str) -> None:
        """Establece notas de orador."""
        text = text.replace("\r\n", "\n").replace("\r", "\n")
        slide = self.get_slide(slide_index)
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = text

    def get_notes(self, slide_index: int) -> str:
        """Obtiene notas de orador."""
        slide = self.get_slide(slide_index)
        return slide.notes_slide.notes_text_frame.text

    # ── Images ─────────────────────────────────────────────────────────────

    def add_image_from_bytes(self, slide_index: int, image_bytes: bytes,
                             left: float = 1.0, top: float = 1.5,
                             width: float = 8.0, height: Optional[float] = None) -> None:
        """Añade imagen desde bytes."""
        slide = self.get_slide(slide_index)
        image_stream = io.BytesIO(image_bytes)
        if height:
            slide.shapes.add_picture(
                image_stream, Inches(left), Inches(top),
                width=Inches(width), height=Inches(height)
            )
        else:
            slide.shapes.add_picture(
                image_stream, Inches(left), Inches(top),
                width=Inches(width)
            )

    def add_image_from_url(self, slide_index: int, url: str, **kwargs) -> None:
        """Descarga y añade imagen desde URL."""
        response = requests.get(url, timeout=15)
        response.raise_for_status()
        self.add_image_from_bytes(slide_index, response.content, **kwargs)

    def add_image_from_path(self, slide_index: int, path: str, **kwargs) -> None:
        """Añade imagen desde archivo local."""
        p = Path(path)
        if not p.exists():
            raise FileNotFoundError(f"Image not found: {path}")
        self.add_image_from_bytes(slide_index, p.read_bytes(), **kwargs)

    def add_ai_image_placeholder(self, slide_index: int, prompt: str,
                                 left: float = 1.0, top: float = 1.5,
                                 width: float = 8.0, height: float = 5.0) -> None:
        """
        Añade un placeholder visual indicando donde irá una imagen IA.
        """
        slide = self.get_slide(slide_index)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(15, 15, 26)
        shape.line.color.rgb = RGBColor(0, 229, 204)

        tf = shape.text_frame
        tf.text = f"🎨 AI IMAGE\n{prompt[:100]}..."
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0, 229, 204)
        p.alignment = PP_ALIGN.CENTER

    # ── Text ───────────────────────────────────────────────────────────────

    def add_text_box(self, slide_index: int, text: str,
                     left: float = 1.0, top: float = 0.5,
                     width: float = 8.0, height: float = 1.5,
                     font_size: int = 24, bold: bool = False,
                     italic: bool = False,
                     color: Tuple[int, int, int] = (255, 255, 255),
                     alignment=PP_ALIGN.LEFT,
                     font_name: str = "Calibri") -> None:
        """Añade caja de texto con estilo."""
        slide = self.get_slide(slide_index)
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = alignment
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = RGBColor(*color)
        run.font.name = font_name
        return txBox

    def add_multiline_text(self, slide_index: int, lines: List[str],
                           left: float = 1.0, top: float = 0.5,
                           width: float = 8.0, height: float = 1.5,
                           font_size: int = 16, bold: bool = False,
                           color: Tuple[int, int, int] = (255, 255, 255),
                           alignment=PP_ALIGN.LEFT,
                           font_name: str = "Calibri",
                           line_spacing: float = 1.5) -> None:
        """Añade texto con múltiples párrafos."""
        slide = self.get_slide(slide_index)
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.alignment = alignment
            p.space_after = Pt(font_size * line_spacing)
            run = p.add_run()
            run.text = line
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(*color)
            run.font.name = font_name

        return txBox

    # ── Shapes ─────────────────────────────────────────────────────────────

    def add_shape(self, slide_index: int, shape_type,
                  left: float, top: float,
                  width: float, height: float,
                  fill_color: Optional[Tuple[int, int, int]] = None,
                  border_color: Optional[Tuple[int, int, int]] = None,
                  border_width: Optional[float] = None,
                  rounded: bool = False) -> Any:
        """Añade una forma con estilo."""
        slide = self.get_slide(slide_index)
        shape_id = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_id, Inches(left), Inches(top),
                                       Inches(width), Inches(height))
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*fill_color)
        else:
            shape.fill.background()

        if border_color and border_width:
            shape.line.color.rgb = RGBColor(*border_color)
            shape.line.width = Pt(border_width)
        else:
            shape.line.fill.background()

        return shape

    def add_line(self, slide_index: int,
                 x1: float, y1: float, x2: float, y2: float,
                 color: Tuple[int, int, int] = (255, 255, 255),
                 width: float = 1.0) -> Any:
        """Añade una línea."""
        slide = self.get_slide(slide_index)
        line = slide.shapes.add_connector(
            MSO_SHAPE.LINE_STRAIGHT_CONNECTOR,
            Inches(x1), Inches(y1), Inches(x2), Inches(y2)
        )
        line.line.color.rgb = RGBColor(*color)
        line.line.width = Pt(width)
        return line

    def add_bar(self, slide_index: int,
                left: float, top: float,
                width: float, height: float,
                color: Tuple[int, int, int]) -> Any:
        """Añade una barra decorativa (rectángulo sin borde)."""
        return self.add_shape(slide_index, MSO_SHAPE.RECTANGLE,
                              left, top, width, height,
                              fill_color=color)

    # ── Background ─────────────────────────────────────────────────────────

    def set_slide_background_color(self, slide_index: int,
                                   r: int, g: int, b: int) -> None:
        """Establece color de fondo sólido."""
        slide = self.get_slide(slide_index)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)

    # ── Info ───────────────────────────────────────────────────────────────

    def summary(self) -> dict:
        if not self.prs:
            return {"slides": 0}
        return {
            "slides": len(self.prs.slides),
            "slide_width_inches": round(self.prs.slide_width / 914400, 2),
            "slide_height_inches": round(self.prs.slide_height / 914400, 2),
        }